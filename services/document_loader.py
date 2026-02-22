"""
Utilities for downloading documents from S3 and extracting their textual contents.
"""
from __future__ import annotations

import io
import os
import tempfile
from dataclasses import dataclass
from html.parser import HTMLParser
from pathlib import Path
from typing import Optional
from botocore.exceptions import ClientError
from docx import Document
from pypdf import PdfReader

from init.aws import get_aws_config, get_s3_client


SUPPORTED_EXTENSIONS = {
    ".txt",
    ".md",
    ".csv",
    ".json",
    ".html",
    ".htm",
    ".pdf",
    ".doc",
    ".docx",
    ".xlsx",
    ".pptx",
}

TEXTUAL_MIME_PREFIXES = ("text/", "application/json", "application/xml")
DOCX_MIME_TYPES = {
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "application/msword",
}
XLSX_MIME_TYPES = {
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    "application/vnd.ms-excel",
}
PPTX_MIME_TYPES = {
    "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    "application/vnd.ms-powerpoint",
}


class _SimpleHTMLTextExtractor(HTMLParser):
    def __init__(self) -> None:
        super().__init__()
        self._parts: list[str] = []

    def handle_data(self, data: str) -> None:
        if data and data.strip():
            self._parts.append(data.strip())

    def text(self) -> str:
        return "\n".join(self._parts).strip()


class DocumentDownloadError(RuntimeError):
    """Raised when the document cannot be downloaded from S3."""


class DocumentParseError(RuntimeError):
    """Raised when the downloaded document cannot be parsed."""


@dataclass(frozen=True, slots=True)
class DocumentContent:
    bucket: str
    key: str
    filename: str
    content_type: Optional[str]
    text: str


class DocumentLoader:
    """
    Fetches documents from S3 and extracts text depending on their type.
    """

    def __init__(self, *, bucket: Optional[str] = None) -> None:
        cfg = get_aws_config()
        self.bucket = bucket or cfg.bucket
        self.client = get_s3_client()

    def fetch_text(self, key: str, *, bucket: Optional[str] = None) -> DocumentContent:
        """
        Download a document from S3 and return its textual contents.
        """
        target_bucket = bucket or self.bucket
        try:
            response = self.client.get_object(Bucket=target_bucket, Key=key)
        except ClientError as exc:
            raise DocumentDownloadError(f"Failed to download {key!r} from S3: {exc}") from exc

        data = response["Body"].read()
        content_type = response.get("ContentType")
        filename = Path(key).name or key
        text = self._extract_text(data, filename, content_type)
        return DocumentContent(bucket=target_bucket, key=key, filename=filename, content_type=content_type, text=text)

    # ------------------------------------------------------------------
    # Text extraction helpers
    # ------------------------------------------------------------------
    def _extract_text(self, data: bytes, filename: str, content_type: Optional[str]) -> str:
        suffix = Path(filename).suffix.lower()
        if suffix and suffix not in SUPPORTED_EXTENSIONS:
            raise DocumentParseError(f"Unsupported file extension: {suffix}")

        if suffix in {".txt", ".md", ".csv", ".json"} or self._is_plain_text(content_type):
            return data.decode("utf-8", errors="replace")

        if suffix in {".html", ".htm"} or content_type in {"text/html", "application/xhtml+xml"}:
            return self._read_html(data)

        if suffix == ".pdf" or content_type == "application/pdf":
            return self._read_pdf(data)

        if suffix in {".docx"} or content_type in DOCX_MIME_TYPES:
            return self._read_docx(data)

        if suffix == ".xlsx" or content_type in XLSX_MIME_TYPES:
            return self._read_xlsx(data)

        if suffix == ".pptx" or content_type in PPTX_MIME_TYPES:
            return self._read_pptx(data)

        if suffix == ".doc":
            return self._read_doc_binary(data)

        raise DocumentParseError("Unknown document format")

    @staticmethod
    def _is_plain_text(content_type: Optional[str]) -> bool:
        if not content_type:
            return False
        return content_type.startswith(TEXTUAL_MIME_PREFIXES)

    @staticmethod
    def _read_pdf(data: bytes) -> str:
        try:
            reader = PdfReader(io.BytesIO(data))
            pages = [page.extract_text() or "" for page in reader.pages]
            return "\n".join(pages).strip()
        except Exception as exc:  # noqa: BLE001
            raise DocumentParseError(f"Failed to parse PDF: {exc}") from exc

    @staticmethod
    def _read_docx(data: bytes) -> str:
        try:
            document = Document(io.BytesIO(data))
            texts = [para.text for para in document.paragraphs]
            return "\n".join(texts).strip()
        except Exception as exc:  # noqa: BLE001
            raise DocumentParseError(f"Failed to parse DOCX: {exc}") from exc

    @staticmethod
    def _read_html(data: bytes) -> str:
        try:
            parser = _SimpleHTMLTextExtractor()
            parser.feed(data.decode("utf-8", errors="replace"))
            parser.close()
            return parser.text()
        except Exception as exc:  # noqa: BLE001
            raise DocumentParseError(f"Failed to parse HTML: {exc}") from exc

    @staticmethod
    def _read_xlsx(data: bytes) -> str:
        try:
            from openpyxl import load_workbook  # type: ignore
        except ImportError as exc:
            raise DocumentParseError(
                "openpyxl package is required to parse .xlsx files."
            ) from exc

        try:
            workbook = load_workbook(filename=io.BytesIO(data), data_only=True, read_only=True)
            lines: list[str] = []
            for sheet in workbook.worksheets:
                lines.append(f"[sheet] {sheet.title}")
                for row in sheet.iter_rows(values_only=True):
                    values = [str(cell).strip() for cell in row if cell is not None and str(cell).strip()]
                    if values:
                        lines.append(" | ".join(values))
            return "\n".join(lines).strip()
        except Exception as exc:  # noqa: BLE001
            raise DocumentParseError(f"Failed to parse XLSX: {exc}") from exc

    @staticmethod
    def _read_pptx(data: bytes) -> str:
        try:
            from pptx import Presentation  # type: ignore
        except ImportError as exc:
            raise DocumentParseError(
                "python-pptx package is required to parse .pptx files."
            ) from exc

        try:
            presentation = Presentation(io.BytesIO(data))
            lines: list[str] = []
            for slide_idx, slide in enumerate(presentation.slides, start=1):
                lines.append(f"[slide] {slide_idx}")
                for shape in slide.shapes:
                    text = getattr(shape, "text", None)
                    if isinstance(text, str) and text.strip():
                        lines.append(text.strip())
            return "\n".join(lines).strip()
        except Exception as exc:  # noqa: BLE001
            raise DocumentParseError(f"Failed to parse PPTX: {exc}") from exc

    @staticmethod
    def _read_doc_binary(data: bytes) -> str:
        try:
            import textract  # type: ignore
        except ImportError as exc:
            raise DocumentParseError(
                "textract package is required to parse .doc files. Install optional dependencies."
            ) from exc

        with tempfile.NamedTemporaryFile(suffix=".doc", delete=False) as tmp:
            tmp.write(data)
            tmp_path = tmp.name

        try:
            extracted = textract.process(tmp_path)
            return extracted.decode("utf-8", errors="replace").strip()
        except Exception as exc:  # noqa: BLE001
            raise DocumentParseError(f"Failed to parse DOC document: {exc}") from exc
        finally:
            try:
                os.remove(tmp_path)
            except OSError:
                pass
